# mcp_server.py
import sys
from mcp.server.fastmcp import FastMCP
from mcp.types import TextContent

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

mcp = FastMCP("PPT Automator")

@mcp.tool()
async def update_slide(
    text: str,
    ppt_path: str = "rectangle_demo.pptx",
    slide_index: int = 1
) -> dict:
    # Load the deck
    prs = Presentation(ppt_path)
    slide = prs.slides[slide_index - 1]

    # Find an existing rectangle
    rect = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE.RECTANGLE:
            rect = shape
            break

    # If none, draw one at 1"×1", 4"×2"
    if rect is None:
        left, top, w, h = Inches(1), Inches(1), Inches(4), Inches(2)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)

    # Put the text in
    rect.text = text

    # Save back to disk
    prs.save(ppt_path)

    return {
        "content": [
            TextContent(
                type="text",
                text=f"✅ Slide {slide_index} updated with '{text}'"
            )
        ]
    }

if __name__ == "__main__":
    # stdio transport for the TS client
    mcp.run(transport="stdio")
